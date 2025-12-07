"""
文档处理器

处理各种类型的文档，包括PDF、Word、Excel等。
"""

import hashlib
import re
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple

import docx
import pandas as pd
import PyPDF2
from pptx import Presentation

from config.logging import get_logger

logger = get_logger(__name__)


class TextCleaner:
    """文本清理器"""

    @staticmethod
    def clean_text(text: str) -> str:
        """
        清理文本

        Args:
            text: 原始文本

        Returns:
            清理后的文本
        """
        if not text:
            return ""

        # 移除多余的空白字符
        text = re.sub(r'\n\s*\n', '\n\n', text)  # 多个空行保留为两个
        text = re.sub(r'\n[ \t]+', '\n', text)  # 移除行尾空格
        text = re.sub(r'[ \t]+', ' ', text)  # 多个空格保留为一个

        # 移除控制字符
        text = ''.join(char for char in text if ord(char) >= 32 or char in '\n\r\t')

        # 移除页眉页脚等模式
        patterns_to_remove = [
            r'第\s*\d+\s*页',
            r'Page\s*\d+',
            r'\f',  # 换页符
        ]

        for pattern in patterns_to_remove:
            text = re.sub(pattern, '', text, flags=re.IGNORECASE)

        return text.strip()

    @staticmethod
    def split_into_paragraphs(text: str, min_length: int = 20) -> List[str]:
        """
        将文本分割成段落

        Args:
            text: 文本
            min_length: 最小段落长度

        Returns:
            段落列表
        """
        # 按双换行分割
        paragraphs = re.split(r'\n\s*\n', text)

        # 过滤短段落
        filtered = [p.strip() for p in paragraphs if len(p.strip()) >= min_length]

        return filtered


class DocumentChunker:
    """文档分块器"""

    def __init__(
        self,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        min_chunk_size: int = 100,
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.min_chunk_size = min_chunk_size

    def chunk_text(
        self,
        text: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> List[Dict[str, Any]]:
        """
        将文本分块

        Args:
            text: 文本
            metadata: 元数据

        Returns:
            文档块列表
        """
        if not text:
            return []

        chunks = []

        # 清理文本
        text = TextCleaner.clean_text(text)

        # 按段落分割
        paragraphs = TextCleaner.split_into_paragraphs(text)

        # 创建文档块
        current_chunk = ""
        current_chunk_paras = []

        for i, para in enumerate(paragraphs):
            # 如果添加当前段落后会超过块大小，则保存当前块
            if len(current_chunk) + len(para) > self.chunk_size and current_chunk:
                chunk = self._create_chunk(
                    current_chunk,
                    current_chunk_paras,
                    chunks,
                    metadata
                )
                chunks.append(chunk)

                # 创建新块（包含重叠部分）
                if len(current_chunk) > self.chunk_overlap:
                    # 从当前块末尾找重叠部分
                    overlap_text = self._find_overlap_text(current_chunk, paragraphs, i)
                    current_chunk = overlap_text
                    current_chunk_paras = []  # 需要重新计算

            # 添加段落到当前块
            current_chunk += para + "\n\n"
            current_chunk_paras.append(para)

        # 添加最后一个块
        if current_chunk and len(current_chunk) >= self.min_chunk_size:
            chunk = self._create_chunk(
                current_chunk,
                current_chunk_paras,
                chunks,
                metadata
            )
            chunks.append(chunk)

        return chunks

    def _find_overlap_text(self, chunk: str, paragraphs: List[str], current_index: int) -> str:
        """找到重叠文本"""
        # 简单实现：从块末尾开始向前找，直到达到重叠大小
        chunk_lines = chunk.split('\n')
        overlap_text = ""
        current_size = 0

        for line in reversed(chunk_lines):
            line = line.strip()
            if not line:
                continue

            if current_size + len(line) <= self.chunk_overlap:
                overlap_text = line + "\n" + overlap_text
                current_size += len(line)
            else:
                break

        return overlap_text.strip()

    def _create_chunk(
        self,
        text: str,
        paragraphs: List[str],
        existing_chunks: List[Dict],
        metadata: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """创建文档块"""
        chunk = {
            "text": text.strip(),
            "paragraphs": paragraphs,
            "chunk_index": len(existing_chunks),
            "chunk_size": len(text),
            "metadata": metadata or {},
        }

        # 添加文本统计信息
        chunk["word_count"] = len(text.split())
        chunk["sentence_count"] = len(re.split(r'[.!?]+', text))

        # 计算文本哈希
        chunk["text_hash"] = hashlib.md5(text.encode()).hexdigest()

        return chunk


class DocumentProcessor:
    """文档处理器"""

    def __init__(self, chunker: Optional[DocumentChunker] = None):
        self.chunker = chunker or DocumentChunker()
        self.supported_formats = {
            '.txt': self._process_text,
            '.md': self._process_text,
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_docx,  # 需要转换
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.pptx': self._process_pptx,
            '.ppt': self._process_pptx,  # 需要转换
        }

    async def process_document(
        self,
        file_path: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        处理文档

        Args:
            file_path: 文件路径
            metadata: 元数据

        Returns:
            处理结果
        """
        try:
            path = Path(file_path)

            if not path.exists():
                return {
                    "success": False,
                    "error": "文件不存在",
                }

            # 检查文件格式
            file_ext = path.suffix.lower()
            if file_ext not in self.supported_formats:
                return {
                    "success": False,
                    "error": f"不支持的文件格式: {file_ext}",
                }

            # 读取并处理文件
            processor_func = self.supported_formats[file_ext]
            text = await processor_func(str(path))

            if not text:
                return {
                    "success": False,
                    "error": "无法从文件中提取文本",
                }

            # 更新元数据
            if not metadata:
                metadata = {}

            metadata.update({
                "file_path": str(path),
                "file_name": path.name,
                "file_ext": file_ext,
                "file_size": path.stat().st_size,
                "processed_at": pd.Timestamp.now().isoformat(),
            })

            # 分块
            chunks = self.chunker.chunk_text(text, metadata)

            # 创建处理结果
            result = {
                "success": True,
                "text": text,
                "chunks": chunks,
                "metadata": metadata,
                "statistics": {
                    "total_length": len(text),
                    "chunk_count": len(chunks),
                    "word_count": len(text.split()),
                    "line_count": len(text.split('\n')),
                },
            }

            logger.info(f"文档处理成功: {file_path}, {len(chunks)} 个块")
            return result

        except Exception as e:
            logger.error(f"处理文档失败 {file_path}: {str(e)}")
            return {
                "success": False,
                "error": str(e),
            }

    async def _process_text(self, file_path: str) -> str:
        """处理纯文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            for encoding in ['gbk', 'gb2312', 'latin-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            raise

    async def _process_pdf(self, file_path: str) -> str:
        """处理PDF文件"""
        text = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text.append(f"=== 页面 {page_num + 1} ===\n{page_text}")
            return '\n\n'.join(text)
        except Exception as e:
            logger.error(f"PDF处理失败: {str(e)}")
            return ""

    async def _process_docx(self, file_path: str) -> str:
        """处理Word文档"""
        try:
            doc = docx.Document(file_path)
            text = []

            # 处理段落
            for para in doc.paragraphs:
                if para.text.strip():
                    text.append(para.text)

            # 处理表格
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    table_text.append(' | '.join(row_text))
                text.append('\n'.join(table_text))

            return '\n\n'.join(text)
        except Exception as e:
            logger.error(f"Word文档处理失败: {str(e)}")
            return ""

    async def _process_excel(self, file_path: str) -> str:
        """处理Excel文件"""
        try:
            # 读取所有工作表
            excel_file = pd.ExcelFile(file_path)
            text = []

            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text.append(f"=== 工作表: {sheet_name} ===")
                text.append(df.to_string(index=False))

            return '\n\n'.join(text)
        except Exception as e:
            logger.error(f"Excel处理失败: {str(e)}")
            return ""

    async def _process_pptx(self, file_path: str) -> str:
        """处理PowerPoint文件"""
        try:
            prs = Presentation(file_path)
            text = []

            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                slide_text.append(f"=== 幻灯片 {slide_num + 1} ===")

                # 处理文本框
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                    # 处理表格
                    if hasattr(shape, "table"):
                        table = shape.table
                        for row in table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            slide_text.append(' | '.join(row_text))

                if slide_text[1:]:  # 除了标题
                    text.append('\n'.join(slide_text))

            return '\n\n'.join(text)
        except Exception as e:
            logger.error(f"PowerPoint处理失败: {str(e)}")
            return ""

    def batch_process(
        self,
        file_paths: List[str],
        metadata_template: Optional[Dict[str, Any]] = None
    ) -> List[Dict[str, Any]]:
        """
        批量处理文档

        Args:
            file_paths: 文件路径列表
            metadata_template: 元数据模板

        Returns:
            处理结果列表
        """
        results = []
        for file_path in file_paths:
            # 创建特定文件的元数据
            metadata = metadata_template.copy() if metadata_template else {}
            metadata["batch_processed"] = True

            # 处理文档
            result = self.process_document(file_path, metadata)
            results.append(result)

        return results

    def update_chunking_strategy(
        self,
        chunk_size: Optional[int] = None,
        chunk_overlap: Optional[int] = None,
        min_chunk_size: Optional[int] = None,
    ):
        """
        更新分块策略

        Args:
            chunk_size: 块大小
            chunk_overlap: 块重叠
            min_chunk_size: 最小块大小
        """
        if chunk_size is not None:
            self.chunker.chunk_size = chunk_size
        if chunk_overlap is not None:
            self.chunker.chunk_overlap = chunk_overlap
        if min_chunk_size is not None:
            self.chunker.min_chunk_size = min_chunk_size

        logger.info(f"更新分块策略: size={self.chunker.chunk_size}, "
                   f"overlap={self.chunker.chunk_overlap}, "
                   f"min={self.chunker.min_chunk_size}")