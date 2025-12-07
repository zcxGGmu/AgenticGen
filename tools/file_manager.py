"""
文件管理器

提供文件上传、下载、管理等功能。
"""

import hashlib
import mimetypes
import os
import shutil
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Any, BinaryIO, Union

from config import settings
from config.logging import get_logger

logger = get_logger(__name__)


class FileManager:
    """文件管理器"""

    def __init__(self, upload_dir: Optional[str] = None):
        self.upload_dir = upload_dir or settings.upload_path
        self.max_file_size = self._parse_size(settings.max_file_size)
        self.allowed_extensions = settings.allowed_extensions

        # 确保上传目录存在
        Path(self.upload_dir).mkdir(parents=True, exist_ok=True)

    def _parse_size(self, size_str: str) -> int:
        """解析文件大小字符串"""
        size_str = size_str.upper()
        if size_str.endswith("KB"):
            return int(size_str[:-2]) * 1024
        elif size_str.endswith("MB"):
            return int(size_str[:-2]) * 1024 * 1024
        elif size_str.endswith("GB"):
            return int(size_str[:-2]) * 1024 * 1024 * 1024
        else:
            return int(size_str)

    def _generate_file_path(self, filename: str, subfolder: Optional[str] = None) -> str:
        """
        生成文件路径

        Args:
            filename: 文件名
            subfolder: 子文件夹

        Returns:
            文件路径
        """
        # 生成唯一文件名
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        name, ext = os.path.splitext(filename)
        unique_filename = f"{name}_{timestamp}{ext}"

        # 添加子文件夹
        if subfolder:
            folder_path = Path(self.upload_dir) / subfolder
            folder_path.mkdir(parents=True, exist_ok=True)
            file_path = folder_path / unique_filename
        else:
            file_path = Path(self.upload_dir) / unique_filename

        return str(file_path)

    def _calculate_file_hash(self, file_path: str) -> str:
        """计算文件哈希值"""
        hash_sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_sha256.update(chunk)
        return hash_sha256.hexdigest()

    async def save_uploaded_file(
        self,
        file_content: bytes,
        filename: str,
        content_type: Optional[str] = None,
        user_id: Optional[str] = None,
        thread_id: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        保存上传的文件

        Args:
            file_content: 文件内容
            filename: 文件名
            content_type: 文件类型
            user_id: 用户ID
            thread_id: 线程ID

        Returns:
            保存结果
        """
        try:
            # 检查文件大小
            file_size = len(file_content)
            if file_size > self.max_file_size:
                return {
                    "success": False,
                    "error": f"文件过大，最大允许: {settings.max_file_size}",
                }

            # 检查文件扩展名
            file_ext = Path(filename).suffix.lower().lstrip('.')
            if file_ext not in self.allowed_extensions:
                return {
                    "success": False,
                    "error": f"不支持的文件类型: {file_ext}",
                }

            # 推断内容类型
            if not content_type:
                content_type, _ = mimetypes.guess_type(filename)

            # 生成文件路径
            user_folder = f"user_{user_id}" if user_id else "anonymous"
            file_path = self._generate_file_path(filename, user_folder)

            # 保存文件
            with open(file_path, "wb") as f:
                f.write(file_content)

            # 计算文件哈希
            file_hash = self._calculate_file_hash(file_path)

            # 构建文件信息
            file_info = {
                "success": True,
                "filename": os.path.basename(file_path),
                "original_filename": filename,
                "file_path": file_path,
                "file_size": file_size,
                "file_type": file_ext,
                "mime_type": content_type,
                "content_hash": file_hash,
                "user_id": user_id,
                "thread_id": thread_id,
                "upload_time": datetime.utcnow().isoformat(),
                "relative_path": os.path.relpath(file_path, self.upload_dir),
            }

            logger.info(f"文件保存成功: {filename} -> {file_path}")
            return file_info

        except Exception as e:
            logger.error(f"保存文件失败: {str(e)}")
            return {
                "success": False,
                "error": f"保存文件失败: {str(e)}",
            }

    def get_file_info(self, file_path: str) -> Optional[Dict[str, Any]]:
        """
        获取文件信息

        Args:
            file_path: 文件路径

        Returns:
            文件信息
        """
        try:
            path = Path(file_path)

            if not path.exists():
                return None

            # 获取基本文件信息
            stat = path.stat()
            file_ext = path.suffix.lower().lstrip('.')
            content_type, _ = mimetypes.guess_type(str(path))

            return {
                "filename": path.name,
                "file_path": str(path),
                "file_size": stat.st_size,
                "file_type": file_ext,
                "mime_type": content_type,
                "created_time": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "modified_time": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "is_file": path.is_file(),
                "is_dir": path.is_dir(),
            }

        except Exception as e:
            logger.error(f"获取文件信息失败: {str(e)}")
            return None

    def list_files(
        self,
        directory: Optional[str] = None,
        user_id: Optional[str] = None,
        recursive: bool = False,
        filter_ext: Optional[List[str]] = None,
    ) -> List[Dict[str, Any]]:
        """
        列出文件

        Args:
            directory: 目录
            user_id: 用户ID
            recursive: 是否递归
            filter_ext: 过滤扩展名

        Returns:
            文件列表
        """
        try:
            # 确定搜索路径
            if user_id:
                search_path = Path(self.upload_dir) / f"user_{user_id}"
            elif directory:
                search_path = Path(self.upload_dir) / directory
            else:
                search_path = Path(self.upload_dir)

            if not search_path.exists():
                return []

            files = []

            # 遍历文件
            pattern = "**/*" if recursive else "*"
            for file_path in search_path.glob(pattern):
                if file_path.is_file():
                    file_ext = file_path.suffix.lower().lstrip('.')

                    # 过滤扩展名
                    if filter_ext and file_ext not in filter_ext:
                        continue

                    file_info = self.get_file_info(str(file_path))
                    if file_info:
                        file_info["relative_path"] = str(file_path.relative_to(self.upload_dir))
                        files.append(file_info)

            return files

        except Exception as e:
            logger.error(f"列出文件失败: {str(e)}")
            return []

    def delete_file(self, file_path: str) -> Dict[str, Any]:
        """
        删除文件

        Args:
            file_path: 文件路径

        Returns:
            删除结果
        """
        try:
            path = Path(file_path)

            if not path.exists():
                return {
                    "success": False,
                    "error": "文件不存在",
                }

            # 安全检查：确保文件在上传目录内
            try:
                path.relative_to(Path(self.upload_dir))
            except ValueError:
                return {
                    "success": False,
                    "error": "无权限删除此文件",
                }

            # 删除文件
            if path.is_file():
                path.unlink()
            elif path.is_dir():
                shutil.rmtree(path)

            logger.info(f"文件删除成功: {file_path}")
            return {
                "success": True,
                "message": "文件删除成功",
            }

        except Exception as e:
            logger.error(f"删除文件失败: {str(e)}")
            return {
                "success": False,
                "error": f"删除文件失败: {str(e)}",
            }

    def move_file(
        self,
        source_path: str,
        target_path: str,
        overwrite: bool = False,
    ) -> Dict[str, Any]:
        """
        移动文件

        Args:
            source_path: 源路径
            target_path: 目标路径
            overwrite: 是否覆盖

        Returns:
            移动结果
        """
        try:
            src = Path(source_path)
            dst = Path(target_path)

            if not src.exists():
                return {
                    "success": False,
                    "error": "源文件不存在",
                }

            # 安全检查
            try:
                src.relative_to(Path(self.upload_dir))
                dst.relative_to(Path(self.upload_dir))
            except ValueError:
                return {
                    "success": False,
                    "error": "无权限操作此文件",
                }

            # 检查目标文件
            if dst.exists() and not overwrite:
                return {
                    "success": False,
                    "error": "目标文件已存在",
                }

            # 创建目标目录
            dst.parent.mkdir(parents=True, exist_ok=True)

            # 移动文件
            shutil.move(str(src), str(dst))

            logger.info(f"文件移动成功: {source_path} -> {target_path}")
            return {
                "success": True,
                "message": "文件移动成功",
                "new_path": str(dst),
            }

        except Exception as e:
            logger.error(f"移动文件失败: {str(e)}")
            return {
                "success": False,
                "error": f"移动文件失败: {str(e)}",
            }

    def get_file_content(self, file_path: str, as_text: bool = True) -> Optional[Union[str, bytes]]:
        """
        读取文件内容

        Args:
            file_path: 文件路径
            as_text: 是否以文本形式读取

        Returns:
            文件内容
        """
        try:
            path = Path(file_path)

            if not path.exists():
                return None

            # 安全检查
            try:
                path.relative_to(Path(self.upload_dir))
            except ValueError:
                logger.error(f"无权限访问文件: {file_path}")
                return None

            # 读取文件
            if as_text:
                with open(path, 'r', encoding='utf-8') as f:
                    return f.read()
            else:
                with open(path, 'rb') as f:
                    return f.read()

        except Exception as e:
            logger.error(f"读取文件失败: {str(e)}")
            return None

    async def extract_text_from_file(self, file_path: str) -> Optional[str]:
        """
        从文件中提取文本

        Args:
            file_path: 文件路径

        Returns:
            提取的文本
        """
        try:
            path = Path(file_path)
            file_ext = path.suffix.lower()

            # 纯文本文件
            if file_ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml']:
                return self.get_file_content(file_path, as_text=True)

            # 需要特殊处理的文件类型
            if file_ext == '.pdf':
                return await self._extract_from_pdf(file_path)
            elif file_ext in ['.docx', '.doc']:
                return await self._extract_from_docx(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return await self._extract_from_excel(file_path)
            elif file_ext == '.pptx':
                return await self._extract_from_pptx(file_path)

            return None

        except Exception as e:
            logger.error(f"提取文本失败: {str(e)}")
            return None

    async def _extract_from_pdf(self, file_path: str) -> Optional[str]:
        """从PDF提取文本"""
        try:
            import PyPDF2
            text = []
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text.append(page.extract_text())
            return '\n'.join(text)
        except ImportError:
            logger.warning("PyPDF2未安装，无法处理PDF文件")
            return None
        except Exception as e:
            logger.error(f"PDF提取失败: {str(e)}")
            return None

    async def _extract_from_docx(self, file_path: str) -> Optional[str]:
        """从Word文档提取文本"""
        try:
            import docx
            doc = docx.Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return '\n'.join(text)
        except ImportError:
            logger.warning("python-docx未安装，无法处理Word文档")
            return None
        except Exception as e:
            logger.error(f"Word文档提取失败: {str(e)}")
            return None

    async def _extract_from_excel(self, file_path: str) -> Optional[str]:
        """从Excel提取文本"""
        try:
            import pandas as pd
            text = []
            # 读取所有工作表
            excel_file = pd.ExcelFile(file_path)
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text.append(f"工作表: {sheet_name}")
                text.append(df.to_string())
            return '\n\n'.join(text)
        except ImportError:
            logger.warning("pandas未安装，无法处理Excel文件")
            return None
        except Exception as e:
            logger.error(f"Excel提取失败: {str(e)}")
            return None

    async def _extract_from_pptx(self, file_path: str) -> Optional[str]:
        """从PowerPoint提取文本"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        except ImportError:
            logger.warning("python-pptx未安装，无法处理PowerPoint文件")
            return None
        except Exception as e:
            logger.error(f"PowerPoint提取失败: {str(e)}")
            return None

    def cleanup_temp_files(self, max_age_hours: int = 24) -> int:
        """
        清理临时文件

        Args:
            max_age_hours: 最大保留时间（小时）

        Returns:
            清理的文件数量
        """
        try:
            cleaned_count = 0
            current_time = datetime.now()
            max_age = max_age_hours * 3600  # 转换为秒

            for file_path in Path(self.upload_dir).rglob("*"):
                if file_path.is_file():
                    # 检查文件年龄
                    file_time = datetime.fromtimestamp(file_path.stat().st_mtime)
                    age_seconds = (current_time - file_time).total_seconds()

                    if age_seconds > max_age:
                        try:
                            file_path.unlink()
                            cleaned_count += 1
                        except Exception as e:
                            logger.error(f"删除临时文件失败: {str(e)}")

            logger.info(f"清理临时文件: {cleaned_count} 个")
            return cleaned_count

        except Exception as e:
            logger.error(f"清理临时文件失败: {str(e)}")
            return 0